import os
import csv
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from tkinter import Tk, filedialog
import json

# Include the following import for CONFIG
CONFIG_FILE_PATH = 'config.json'

try:
    with open(CONFIG_FILE_PATH, 'r') as config_file:
        CONFIG = json.load(config_file)
except FileNotFoundError:
    CONFIG = {
        'CSV_FILE_NAME': 'urls.csv',
        'OUTPUT_IMAGES_FOLDER': 'output_images',
        'TIMEOUT': 10,
        'RETRIES': 2
    }

def create_textbox(slide, left, top, width, height):
    """
    Create and return a text frame inside a textbox on the slide.
    """
    textbox = slide.shapes.add_textbox(left, top, width, height)
    return textbox.text_frame

def customize_paragraph(p, font_size=7, alignment=2):
    """
    Customize the paragraph with the specified font size and alignment.
    """
    p.font.size = Pt(font_size)
    p.alignment = alignment

def add_footer_to_slide(slide, urls, presentation):
    """
    Add a footer to the slide containing the provided URLs.
    """
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    left = Inches(0.5)
    top = slide_height - Inches(0.5)
    width = slide_width - Inches(1.0)
    height = Inches(0.5)

    text_frame = create_textbox(slide, left, top, width, height)
    
    p = text_frame.add_paragraph()
    for url in urls:
        p.text += url + '\n'
    
    customize_paragraph(p)

def read_csv(csv_file):
    """
    Read data from the CSV file and organize it by slide number.
    """
    data = {}
    with open(csv_file, newline='', encoding='utf-8') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            if row and len(row) >= 2:
                image_path, url = row[0], row[1]

                match = re.search(r'slide_(\d+)_image_\d+', image_path)
                if match:
                    slide_number = int(match.group(1))
                    data.setdefault(slide_number, []).append(url)
    return data

def add_footer_to_slides(input_pptx, output_folder, csv_file):
    """
    Add footers containing URLs to slides based on the CSV data.
    Save the modified presentation to the output folder.
    """
    presentation = Presentation(input_pptx)
    data = read_csv(csv_file)

    # Update the following line to use the parent directory of dynamic_directory
    parent_directory = os.path.dirname(dynamic_directory)
    output_folder = os.path.join(parent_directory, "output_pptx")

    for i, slide in enumerate(presentation.slides):
        if i + 1 in data:
            urls = data[i + 1]
            add_footer_to_slide(slide, urls, presentation)

    os.makedirs(output_folder, exist_ok=True)

    base_name = os.path.splitext(os.path.basename(input_pptx))[0]
    output_path = os.path.join(output_folder, f"{base_name}_modified.pptx")
    presentation.save(output_path)

if __name__ == "__main__":
    dynamic_directory = CONFIG.get('OUTPUT_IMAGES_FOLDER', 'default_dynamic_directory')

    # Define parent_directory based on dynamic_directory
    parent_directory = os.path.dirname(dynamic_directory)

    try:
        # Detect PowerPoint file dynamically in the "input_pptx" folder
        input_pptx_folder = os.path.join(parent_directory, "input_pptx")
        pptx_files = [file for file in os.listdir(input_pptx_folder) if file.endswith('.pptx')]
        
        if not pptx_files:
            print("Error: No PowerPoint files found in the 'input_pptx' folder.")
        else:
            # Use the first found PowerPoint file
            input_pptx = os.path.join(input_pptx_folder, pptx_files[0])

            # Use the correct path for csv_file
            csv_file = os.path.join(parent_directory, CONFIG['CSV_FILE_NAME'])

            add_footer_to_slides(input_pptx, dynamic_directory, csv_file)
    except Exception as e:
        print(f"Error in Script 3: {e}")
