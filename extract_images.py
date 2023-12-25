import os
import logging
from pptx import Presentation
from PIL import Image, ImageFile
from io import BytesIO
from tkinter import filedialog, Tk
import json
from datetime import datetime
import shutil

ImageFile.LOAD_TRUNCATED_IMAGES = True

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def convert_to_rgb(image):
    if image.mode in ("RGBA", "P"):
        return image.convert("RGB")
    return image

def process_image(image, output_folder, slide_number, image_count, template):
    try:
        image_data = image.blob
        with BytesIO(image_data) as image_stream:
            img = Image.open(image_stream)
            img = convert_to_rgb(img)

            img_path = os.path.join(output_folder, f"{template}__fin_slide_{slide_number}_image_{image_count + 1}.jpg")
            img.save(img_path, "JPEG")

            return True
    except Exception as e:
        logger.error(f"Error processing image on slide {slide_number}, shape {image_count + 1}: {e}")
        return False

def extract_and_process_images(presentation, output_folder, template):
    image_count = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "image"):
                if process_image(shape.image, output_folder, slide_number, image_count, template):
                    image_count += 1

    return image_count

def extract_images_from_pptx(input_pptx, base_output_folder, template):
    try:
        # Load the PowerPoint presentation
        presentation = Presentation(input_pptx)

        # Create a new directory using the name of the PowerPoint file (without extension)
        base_directory = os.path.splitext(os.path.basename(input_pptx))[0]
        
        # Dynamically update the output folder based on the current date and time
        current_time_str = datetime.now().strftime("%Y%m%d_%H%M%S")
        dynamic_output_folder = os.path.join(os.getcwd(), f"{base_directory}_{current_time_str}")

        # Create the output folder if it doesn't exist
        os.makedirs(dynamic_output_folder, exist_ok=True)

        # Inside the new directory, create "output_images" and "input_pptx" folders
        output_images_folder = os.path.join(dynamic_output_folder, "output_images")
        os.makedirs(output_images_folder, exist_ok=True)

        input_pptx_folder = os.path.join(dynamic_output_folder, "input_pptx")
        os.makedirs(input_pptx_folder, exist_ok=True)

        # Copy the PowerPoint file to the "input_pptx" folder
        new_pptx_path = os.path.join(input_pptx_folder, f"{base_directory}.pptx")
        shutil.copy(input_pptx, new_pptx_path)

        # Extract and process images
        image_count = extract_and_process_images(presentation, output_images_folder, template)

        # Update the configuration file with the new output folder
        update_config(output_images_folder)

        logger.info(f"{image_count} images extracted from the PowerPoint presentation.")
        return dynamic_output_folder  # Return the path to the new directory
    except Exception as e:
        logger.error(f"Error extracting images from PowerPoint: {e}")

def update_config(new_output_folder):
    # Load the existing configuration from the file
    config_file_path = 'config.json'
    try:
        with open(config_file_path, 'r') as config_file:
            config = json.load(config_file)
    except FileNotFoundError:
        # Provide default values if the configuration file is missing
        config = {
            'CSV_FILE_NAME': 'urls.csv',
            'OUTPUT_IMAGES_FOLDER': '',
            'TIMEOUT': 10,
            'RETRIES': 2
        }

    # Update the OUTPUT_IMAGES_FOLDER value in the configuration
    config['OUTPUT_IMAGES_FOLDER'] = new_output_folder

    # Write the updated configuration back to the file
    with open(config_file_path, 'w') as config_file:
        json.dump(config, config_file, indent=2)

def choose_file():
    Tk().withdraw()  # Hide the main window
    file_path = filedialog.askopenfilename(title="Select PowerPoint file", filetypes=[("PowerPoint files", "*.pptx")])
    return file_path

def main():
    try:
        # Choose the input PowerPoint file interactively
        input_pptx = choose_file()

        # Extract images from PowerPoint and mark slide numbers
        template = os.path.splitext(os.path.basename(input_pptx))[0]

        # Extract images and update configuration
        dynamic_output_folder = extract_images_from_pptx(input_pptx, None, template)

        logger.info(f"Images extracted. Output directory: {dynamic_output_folder}")
    except KeyboardInterrupt:
        logger.info("Script execution interrupted by the user.")

if __name__ == "__main__":
    main()
